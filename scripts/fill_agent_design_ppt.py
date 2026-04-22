"""Generate a 4-slide Agent Design PPT by filling the template.

Each output slide = 1 key agent. The template's 4 "Agent" boxes are
reused as the 4 rubric facets: Purpose / I&O / Reasoning+Memory+Tools /
Interaction.

Run:
    python scripts/fill_agent_design_ppt.py
"""
from __future__ import annotations

import copy
from datetime import datetime
from pathlib import Path

from pptx import Presentation


TEMPLATE = Path(r"c:/Users/Zyan/Downloads/演示文稿1.pptx")
_DEFAULT_OUT = Path(r"c:/Users/Zyan/Downloads/演示文稿1_agent_design_filled.pptx")


def _resolve_output() -> Path:
    """If the default output is locked (open in PowerPoint), add a timestamp."""
    if not _DEFAULT_OUT.exists():
        return _DEFAULT_OUT
    try:
        with _DEFAULT_OUT.open("ab"):
            return _DEFAULT_OUT
    except PermissionError:
        ts = datetime.now().strftime("%Y%m%d_%H%M%S")
        return _DEFAULT_OUT.with_name(f"{_DEFAULT_OUT.stem}_{ts}.pptx")


OUTPUT = _resolve_output()


AGENTS = [
    {
        "slide_title": "3. Agent Design  ·  Intent Profile Agent",
        "subtitle": "Turns a free-form user message into structured travel intent that every downstream agent relies on.",
        "row1_title": "Purpose & Responsibilities",
        "row1_body": (
            "Extract a clean, structured user profile (origin, destination, dates, budget, party, "
            "hard constraints and soft preferences) and generate targeted search queries. "
            "It is the only agent that interprets the raw user text into planning parameters."
        ),
        "row2_title": "Input  &  Output",
        "row2_body": (
            "Input:  sanitised user message + prior state fields (messages, prior intent).\n"
            "Output:  intent_profile_output, hard_constraints, soft_preferences, search_queries."
        ),
        "row3_title": "Reasoning · Memory · Tools",
        "row3_body": (
            "Model:  OPENAI_MODEL (gpt-4.1-nano) via LLM reasoning with a strict JSON schema.\n"
            "Memory:  short-term only — reads/writes the shared LangGraph State.\n"
            "Tools:   none (pure LLM structuring, no external API calls)."
        ),
        "row4_title": "Interaction with Other Agents",
        "row4_body": (
            "Receives from:  Input Guard Agent (sanitised_input).\n"
            "Sends to:       Research Agent (search_queries) and Planner Agent (constraints / preferences)."
        ),
    },
    {
        "slide_title": "3. Agent Design  ·  Research Agent",
        "subtitle": "Collects real-world flight, hotel, weather and attraction data so the planner can ground its itinerary.",
        "row1_title": "Purpose & Responsibilities",
        "row1_body": (
            "Retrieve and consolidate external travel data based on the intent queries: flights, hotels, "
            "attractions and weather. It is the only agent allowed to hit external search APIs for data grounding."
        ),
        "row2_title": "Input  &  Output",
        "row2_body": (
            "Input:  search_queries + hard_constraints from Intent Profile Agent.\n"
            "Output:  research (raw evidence) and inventory (filtered, de-duplicated candidates) in State."
        ),
        "row3_title": "Reasoning · Memory · Tools",
        "row3_body": (
            "Model:  OPENAI_MODEL (gpt-4.1-nano) for query normalisation + rule-based filtering / de-duplication.\n"
            "Memory:  short-term (State); results cached per run so the planner can re-read without re-calling APIs.\n"
            "Tools:   search_flights, search_hotels, search_weather, google_search, web_search."
        ),
        "row4_title": "Interaction with Other Agents",
        "row4_body": (
            "Receives from:  Intent Profile Agent.\n"
            "Sends to:       Planner Agent (inventory / research)  →  which uses it to build candidate itineraries."
        ),
    },
    {
        "slide_title": "3. Agent Design  ·  Planner Agent",
        "subtitle": "Turns structured intent and research evidence into multiple candidate day-by-day itineraries.",
        "row1_title": "Purpose & Responsibilities",
        "row1_body": (
            "Generate 2–3 diverse, budget-aware itinerary options and, when the Debate Agent returns a critique, "
            "revise the plan. It is the only agent that writes itineraries; other agents only evaluate or explain."
        ),
        "row2_title": "Input  &  Output",
        "row2_body": (
            "Input:  intent_profile_output, inventory, and (on revision) critique from Debate Agent.\n"
            "Output:  itineraries, planner_decision_trace; resets is_valid = None to trigger a fresh debate round."
        ),
        "row3_title": "Reasoning · Memory · Tools",
        "row3_body": (
            "Model:  PLANNER_MODEL (gpt-4.1) — LLM reasoning with structured prompts for draft + revise modes.\n"
            "Memory:  short-term State + long-term DB (plan_id persisted so a plan can be revisited / replanned).\n"
            "Tools:   same search tools as Research Agent, used sparingly to fill gaps while planning."
        ),
        "row4_title": "Interaction with Other Agents",
        "row4_body": (
            "Receives from:  Research Agent (inventory), Debate Agent (critique on revision), Replanner (user feedback).\n"
            "Sends to:       Debate Agent (new itineraries)  →  then Explainability and Output Guard."
        ),
    },
    {
        "slide_title": "3. Agent Design  ·  Debate Agent",
        "subtitle": "Critiques the planner's itineraries and drives an iterative improvement loop up to three rounds.",
        "row1_title": "Purpose & Responsibilities",
        "row1_body": (
            "Evaluate itineraries on bias / fairness, logistics, preference alignment and option diversity, then "
            "decide continue (send critique back to Planner) or decide (accept). It owns the loop termination logic."
        ),
        "row2_title": "Input  &  Output",
        "row2_body": (
            "Input:  itineraries, intent_profile_output, debate_count, debate_history.\n"
            "Output:  critique, is_valid (True/False), round_decision, debate_verdict, final_itineraries when accepted."
        ),
        "row3_title": "Reasoning · Memory · Tools",
        "row3_body": (
            "Model:  DEBATE_MODEL (gpt-4.1) for per-round critique + JUDGE_MODEL (gpt-5-mini) for final verdict.\n"
            "Memory:  short-term debate_history in State + long-term DB (debate_verdict persisted with the plan).\n"
            "Tools:   none — pure LLM reasoning over the planner's output."
        ),
        "row4_title": "Interaction with Other Agents",
        "row4_body": (
            "Receives from:  Planner Agent.\n"
            "Sends to:       Planner Agent (critique, while debate_count < MAX_DEBATE_ROUNDS=3) or Explainability Agent (accept)."
        ),
    },
]


def set_text_preserve_format(shape, new_text: str) -> None:
    """Replace a shape's text, keeping the template's font/colour."""
    tf = shape.text_frame
    first_para = tf.paragraphs[0]
    first_rpr = None
    if first_para.runs:
        rpr = first_para.runs[0]._r.find(
            "{http://schemas.openxmlformats.org/drawingml/2006/main}rPr"
        )
        if rpr is not None:
            first_rpr = copy.deepcopy(rpr)

    tf.clear()
    for i, line in enumerate(new_text.split("\n")):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        run = p.add_run()
        run.text = line
        if first_rpr is not None:
            r_el = run._r
            existing = r_el.find(
                "{http://schemas.openxmlformats.org/drawingml/2006/main}rPr"
            )
            if existing is not None:
                r_el.remove(existing)
            r_el.insert(0, copy.deepcopy(first_rpr))


def clone_slide(prs, src_slide):
    """Append a deep copy of src_slide to prs."""
    new_slide = prs.slides.add_slide(src_slide.slide_layout)
    for shp in list(new_slide.shapes):
        shp.element.getparent().remove(shp.element)
    for shp in src_slide.shapes:
        new_slide.shapes._spTree.insert_element_before(
            copy.deepcopy(shp.element), "p:extLst"
        )
    return new_slide


def main() -> None:
    prs = Presentation(str(TEMPLATE))
    src = prs.slides[0]

    slides = [src] + [clone_slide(prs, src) for _ in range(3)]

    field_map = [
        ("Text 3", "Text 4"),
        ("Text 8", "Text 9"),
        ("Text 13", "Text 14"),
        ("Text 18", "Text 19"),
    ]
    detail_shapes = [
        "Text 5", "Text 6",
        "Text 10", "Text 11",
        "Text 15", "Text 16",
        "Text 20", "Text 21",
    ]

    for slide, agent in zip(slides, AGENTS):
        by_name = {s.name: s for s in slide.shapes if s.has_text_frame}

        if "Text 0" in by_name:
            set_text_preserve_format(by_name["Text 0"], agent["slide_title"])
        if "Text 1" in by_name:
            set_text_preserve_format(by_name["Text 1"], agent["subtitle"])

        rows = [
            (agent["row1_title"], agent["row1_body"]),
            (agent["row2_title"], agent["row2_body"]),
            (agent["row3_title"], agent["row3_body"]),
            (agent["row4_title"], agent["row4_body"]),
        ]
        for (t_name, b_name), (t_text, b_text) in zip(field_map, rows):
            if t_name in by_name:
                set_text_preserve_format(by_name[t_name], t_text)
            if b_name in by_name:
                set_text_preserve_format(by_name[b_name], b_text)

        for d in detail_shapes:
            if d in by_name:
                set_text_preserve_format(by_name[d], "")

    prs.save(str(OUTPUT))
    print(f"Saved -> {OUTPUT}")


if __name__ == "__main__":
    main()
